from pathlib import Path
from typing import Any

from automation.engine.context import AutomationContext
from automation.engine.rollback import RollbackManager
from automation.engine.artifacts import artifact_manager


class OfficeEngine:
    """Office document automation: Word, Excel, PowerPoint, PDF."""

    def create_document(self, params: dict, ctx: AutomationContext, rollback: RollbackManager) -> dict:
        doc_type = params.get("type", "word")
        content = params.get("content", "")
        name = params.get("name", "document")

        if doc_type == "word":
            return self._create_word(name, content, ctx)
        elif doc_type == "excel":
            return self._create_excel(name, content, ctx)
        elif doc_type == "powerpoint":
            return self._create_ppt(name, content, ctx)
        elif doc_type == "pdf":
            return self._create_pdf(name, content, ctx)
        return {"status": "error", "error": f"Unknown document type: {doc_type}"}

    def _create_word(self, name: str, content: str, ctx: AutomationContext) -> dict:
        try:
            from docx import Document
            doc = Document()
            for line in content.split("\n"):
                doc.add_paragraph(line)
            import io
            buf = io.BytesIO()
            doc.save(buf)
            artifact = artifact_manager.save_file(name, buf.getvalue(), ctx.automation_id, "docx")
            return {"status": "ok", "artifact_id": artifact.id, "path": artifact.path}
        except ImportError:
            return {"status": "error", "error": "python-docx not installed"}

    def _create_excel(self, name: str, content: str, ctx: AutomationContext) -> dict:
        try:
            from openpyxl import Workbook
            wb = Workbook()
            ws = wb.active
            for i, line in enumerate(content.split("\n"), 1):
                for j, cell in enumerate(line.split("\t"), 1):
                    ws.cell(row=i, column=j, value=cell)
            import io
            buf = io.BytesIO()
            wb.save(buf)
            artifact = artifact_manager.save_file(name, buf.getvalue(), ctx.automation_id, "xlsx")
            return {"status": "ok", "artifact_id": artifact.id, "path": artifact.path}
        except ImportError:
            return {"status": "error", "error": "openpyxl not installed"}

    def _create_ppt(self, name: str, content: str, ctx: AutomationContext) -> dict:
        try:
            from pptx import Presentation
            prs = Presentation()
            for line in content.split("\n"):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = line
            import io
            buf = io.BytesIO()
            prs.save(buf)
            artifact = artifact_manager.save_file(name, buf.getvalue(), ctx.automation_id, "pptx")
            return {"status": "ok", "artifact_id": artifact.id, "path": artifact.path}
        except ImportError:
            return {"status": "error", "error": "python-pptx not installed"}

    def _create_pdf(self, name: str, content: str, ctx: AutomationContext) -> dict:
        try:
            from fpdf import FPDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Helvetica", size=12)
            for line in content.split("\n"):
                pdf.cell(0, 10, text=line, new_x="LMARGIN", new_y="NEXT")
            import io
            buf = io.BytesIO()
            pdf.output(buf)
            artifact = artifact_manager.save_file(name, buf.getvalue(), ctx.automation_id, "pdf")
            return {"status": "ok", "artifact_id": artifact.id, "path": artifact.path}
        except ImportError:
            return {"status": "error", "error": "fpdf2 not installed"}

    def read_document(self, params: dict, ctx: AutomationContext, rollback: RollbackManager) -> dict:
        path = params.get("path", "")
        ext = Path(path).suffix.lower()

        if ext == ".docx":
            try:
                from docx import Document
                doc = Document(path)
                text = "\n".join(p.text for p in doc.paragraphs)
                return {"status": "ok", "content": text[:5000], "length": len(text)}
            except ImportError:
                return {"status": "error", "error": "python-docx not installed"}
        elif ext == ".xlsx":
            try:
                from openpyxl import load_workbook
                wb = load_workbook(path, read_only=True)
                ws = wb.active
                rows = [[cell.value for cell in row] for row in ws.iter_rows(max_row=50)]
                return {"status": "ok", "rows": rows, "row_count": ws.max_row}
            except ImportError:
                return {"status": "error", "error": "openpyxl not installed"}
        elif ext == ".pdf":
            try:
                import fitz
                doc = fitz.open(path)
                text = "\n".join(page.get_text() for page in doc)
                return {"status": "ok", "content": text[:5000], "pages": len(doc)}
            except ImportError:
                return {"status": "error", "error": "PyMuPDF not installed"}

        return {"status": "error", "error": f"Unsupported format: {ext}"}


office_engine = OfficeEngine()
